"""Brand a Markdown file into a MobiCova .docx (teal headings + logo banner),
and build a one-slide budget-summary .pptx. Reusable for the doc kit.

Usage:
  python scripts/brand_docs.py md2docx INPUT.md OUTPUT.docx "Document Title"
  python scripts/brand_docs.py budget-slide OUTPUT.pptx
"""
import sys, subprocess, os
from docx import Document
from docx.shared import RGBColor, Pt, Inches, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO = os.path.join(ROOT, "client", "public", "images", "logo.png")
PANDOC = r"C:\Users\Dee\AppData\Local\Pandoc\pandoc.exe"
TEAL = RGBColor(0x0A, 0x7B, 0x7B)
DARK = RGBColor(0x0E, 0x2A, 0x2A)


def add_logo_banner(doc, title):
    # Fresh centered logo paragraph + teal title paragraph, both moved to the top.
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    try:
        run.add_picture(LOGO, width=Inches(2.4))
    except Exception:
        run.add_text("MobiCova Health")

    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = t.add_run(title)
    tr.font.size = Pt(20)
    tr.font.bold = True
    tr.font.color.rgb = TEAL
    tr.font.name = "Calibri"

    # Move the banner block to the very top (insert title last so it ends up under the logo).
    body = doc.element.body
    for para in [t._p, p._p]:
        body.remove(para)
        body.insert(0, para)


def recolor_headings(doc):
    for para in doc.paragraphs:
        if para.style.name and para.style.name.startswith("Heading"):
            for run in para.runs:
                run.font.color.rgb = TEAL
                run.font.name = "Calibri"


def md2docx(md, out, title):
    tmp = out + ".tmp.docx"
    subprocess.run([PANDOC, md, "-o", tmp], check=True)
    doc = Document(tmp)
    # base font
    try:
        doc.styles["Normal"].font.name = "Calibri"
        doc.styles["Normal"].font.size = Pt(10.5)
    except Exception:
        pass
    # Drop the markdown's own leading H1 (the banner supplies the title).
    for para in doc.paragraphs:
        if para.style.name == "Title" or para.style.name == "Heading 1":
            para._p.getparent().remove(para._p)
            break
    recolor_headings(doc)
    add_logo_banner(doc, title)
    doc.save(out)
    os.remove(tmp)
    print("wrote", out)


def budget_slide(out):
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt, Emu as PEmu
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.text import PP_ALIGN

    TEALP = PRGB(0x0A, 0x7B, 0x7B)
    DARKP = PRGB(0x0E, 0x2A, 0x2A)
    GREY = PRGB(0x55, 0x66, 0x66)
    WHITE = PRGB(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # dark header band
    band = slide.shapes.add_shape(1, 0, 0, prs.slide_width, PInches(1.5))
    band.fill.solid(); band.fill.fore_color.rgb = DARKP; band.line.fill.background()
    try:
        slide.shapes.add_picture(LOGO, PInches(0.5), PInches(0.42), height=PInches(0.66))
    except Exception:
        pass
    tb = slide.shapes.add_textbox(PInches(4.7), PInches(0.45), PInches(8.2), PInches(0.7))
    tf = tb.text_frame; tf.text = "Production Platform — Year-One Budget (50k members)"
    r = tf.paragraphs[0].runs[0]; r.font.size = PPt(22); r.font.bold = True; r.font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    rows = [
        ("Bucket", "Year-one (USD)"),
        ("Infra build-out / setup (Tier B, AWS af-south-1)", "$30k - $70k"),
        ("Infra run-rate  (~$2k-$3k/mo x 12)", "$24k - $36k"),
        ("SOC 2 + ISO 27001 + penetration test", "$40k - $90k"),
        ("DevOps / SRE  (~$2k-$6k/mo x 12)", "$24k - $72k"),
        ("SUBTOTAL  (excl. variable comms & payments)", "$118k - $268k"),
        ("Variable comms (WhatsApp/SMS/USSD) + AI", "$18k - $130k+  (usage)"),
    ]
    top = PInches(1.9); left = PInches(0.6)
    w = prs.slide_width - PInches(1.2); rh = PInches(0.62)
    tbl = slide.shapes.add_table(len(rows), 2, left, top, w, rh*len(rows)).table
    tbl.columns[0].width = PInches(8.6); tbl.columns[1].width = PInches(3.53)
    for ri, (a, b) in enumerate(rows):
        for ci, val in enumerate((a, b)):
            cell = tbl.cell(ri, ci)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0]; run.font.size = PPt(13); run.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEALP
                run.font.color.rgb = WHITE; run.font.bold = True
            elif "SUBTOTAL" in a:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRGB(0xE6, 0xF2, 0xF2)
                run.font.bold = True; run.font.color.rgb = DARKP
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
                run.font.color.rgb = PRGB(0x22, 0x33, 0x33)
            if ci == 1:
                para.alignment = PP_ALIGN.RIGHT

    note = slide.shapes.add_textbox(PInches(0.6), PInches(6.5), w, PInches(0.8))
    nf = note.text_frame; nf.word_wrap = True
    nf.text = ("Planning-grade estimates. At 50k, servers are the small part — setup, compliance, "
               "people and per-message comms drive the budget. Destination: AWS af-south-1 (African "
               "residency + SOC 2 fit). Comms line is highly sensitive to USSD billing model & WhatsApp template mix.")
    nr = nf.paragraphs[0].runs[0]; nr.font.size = PPt(10.5); nr.font.italic = True; nr.font.color.rgb = GREY

    prs.save(out)
    print("wrote", out)


def language_budget_slide(out):
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    TEALP = PRGB(0x0A, 0x7B, 0x7B)
    TEAL2 = PRGB(0x13, 0x9B, 0x9B)
    DARKP = PRGB(0x0E, 0x2A, 0x2A)
    GREY = PRGB(0x55, 0x66, 0x66)
    ORANGE = PRGB(0xE2, 0x7D, 0x1E)
    WHITE = PRGB(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    def header(slide, title):
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PInches(1.5))
        band.fill.solid(); band.fill.fore_color.rgb = DARKP; band.line.fill.background()
        try:
            slide.shapes.add_picture(LOGO, PInches(0.5), PInches(0.42), height=PInches(0.66))
        except Exception:
            pass
        tb = slide.shapes.add_textbox(PInches(4.2), PInches(0.30), PInches(8.7), PInches(0.6))
        tf = tb.text_frame; tf.word_wrap = True
        tf.text = title
        r = tf.paragraphs[0].runs[0]; r.font.size = PPt(21); r.font.bold = True; r.font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        sb = slide.shapes.add_textbox(PInches(4.2), PInches(0.95), PInches(8.7), PInches(0.4))
        sf = sb.text_frame; sf.text = "Prepared for AXA Mansard"
        sr = sf.paragraphs[0].runs[0]; sr.font.size = PPt(13); sr.font.italic = True; sr.font.color.rgb = PRGB(0xBF, 0xD8, 0xD8)
        sf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    header(slide, "Nigerian-Language Support — Investment Summary")

    rows = [
        ("Cost area", "Indicative (all 4 languages)"),
        ("Translation - general (UI + USSD/WhatsApp + email/SMS)", "₦0.8M - 1.4M"),
        ("Translation - specialist (safety / legal / consent)", "₦0.5M - 0.9M"),
        ("Native-speaker clinician review & sign-off", "₦2.0M - 3.0M"),
        ("In-language safety vocabularies", "₦0.3M - 0.5M"),
        ("Legal review of consent / disclaimers", "₦0.6M - 1.2M"),
        ("Linguistic QA & testing", "₦0.2M - 0.4M"),
        ("ONE-OFF TOTAL (external, ~$2.8k - 4.6k)", "₦4.4M - 7.4M"),
        ("Recurring (ongoing)", "₦0.6M - 1.4M / year"),
    ]
    top = PInches(1.85); left = PInches(0.6)
    w = prs.slide_width - PInches(1.2); rh = PInches(0.52)
    tbl = slide.shapes.add_table(len(rows), 2, left, top, w, rh*len(rows)).table
    tbl.columns[0].width = PInches(8.6); tbl.columns[1].width = PInches(3.53)
    for ri, (a, b) in enumerate(rows):
        for ci, val in enumerate((a, b)):
            cell = tbl.cell(ri, ci)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0]; run.font.size = PPt(13); run.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEALP
                run.font.color.rgb = WHITE; run.font.bold = True
            elif "TOTAL" in a:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRGB(0xE6, 0xF2, 0xF2)
                run.font.bold = True; run.font.color.rgb = DARKP
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
                run.font.color.rgb = PRGB(0x22, 0x33, 0x33)
            if ci == 1:
                para.alignment = PP_ALIGN.RIGHT

    note = slide.shapes.add_textbox(PInches(0.6), PInches(6.7), w, PInches(0.7))
    nf = note.text_frame; nf.word_wrap = True
    nf.text = ("Indicative planning ranges, not a quote. Cost is driven by clinical & legal review, not software. "
               "Rollout: Pidgin -> Hausa -> Yoruba -> Igbo; first language live ~6-8 weeks from go-ahead, all four ~4-5 months.")
    nr = nf.paragraphs[0].runs[0]; nr.font.size = PPt(10.5); nr.font.italic = True; nr.font.color.rgb = GREY

    # ---- Slide 2: indicative delivery timeline (Gantt) ----
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    header(s2, "Indicative Delivery Timeline")

    WEEK0_IN = 3.5          # left edge of the timeline area (inches)
    TOTAL_WK = 20
    SPAN_IN = 9.2           # width of the 20-week span
    wpx = SPAN_IN / TOTAL_WK
    def wx(week):
        return PInches(WEEK0_IN + week * wpx)

    phases = [
        ("Phase 0 — Foundation", 0, 3, "Wk 1–3", GREY),
        ("Phase 1 — Nigerian Pidgin", 3, 8, "Wk 3–8", TEALP),
        ("Phase 2 — Hausa", 8, 12, "Wk 8–12", TEAL2),
        ("Phase 3 — Yoruba", 12, 16, "Wk 12–16", TEALP),
        ("Phase 4 — Igbo", 16, 20, "Wk 16–20", TEAL2),
    ]
    row_top0 = 2.30
    pitch = 0.72
    bar_h = 0.46

    # Milestone vertical guide lines (orange) at "first language live" (wk 8) and "all four" (wk 20)
    axis_top = 2.05
    axis_bottom = row_top0 + len(phases) * pitch + 0.05
    for wk in (8, 20):
        ln = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, wx(wk), PInches(axis_top),
                                 PInches(0.022), PInches(axis_bottom - axis_top))
        ln.fill.solid(); ln.fill.fore_color.rgb = ORANGE; ln.line.fill.background()
    for wk, label in ((8, "Pidgin live ~wk 8"), (20, "All four ~wk 20")):
        ml = s2.shapes.add_textbox(PInches(WEEK0_IN + wk * wpx - 1.6), PInches(1.62), PInches(1.7), PInches(0.35))
        mf = ml.text_frame; mf.word_wrap = True; mf.text = label
        mr = mf.paragraphs[0].runs[0]; mr.font.size = PPt(10.5); mr.font.bold = True; mr.font.color.rgb = ORANGE
        mf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Phase rows: left label + bar
    for i, (name, start, end, wklabel, color) in enumerate(phases):
        rt = row_top0 + i * pitch
        lbl = s2.shapes.add_textbox(PInches(0.5), PInches(rt), PInches(2.85), PInches(bar_h))
        lf = lbl.text_frame; lf.word_wrap = True; lf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lf.text = name
        lr = lf.paragraphs[0].runs[0]; lr.font.size = PPt(12.5); lr.font.bold = True; lr.font.color.rgb = DARKP
        lf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        bar = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, wx(start), PInches(rt),
                                  PInches((end - start) * wpx), PInches(bar_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        bf = bar.text_frame; bf.word_wrap = True
        bf.text = wklabel
        br = bf.paragraphs[0].runs[0]; br.font.size = PPt(11); br.font.bold = True; br.font.color.rgb = WHITE
        bf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Week axis line + tick labels
    ax_y = axis_bottom + 0.06
    axis = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, wx(0), PInches(ax_y), PInches(SPAN_IN), PInches(0.018))
    axis.fill.solid(); axis.fill.fore_color.rgb = PRGB(0xC8, 0xD4, 0xD4); axis.line.fill.background()
    for wk in (0, 4, 8, 12, 16, 20):
        tk = s2.shapes.add_textbox(PInches(WEEK0_IN + wk * wpx - 0.5), PInches(ax_y + 0.06), PInches(1.0), PInches(0.3))
        tf2 = tk.text_frame; tf2.text = f"Wk {wk}"
        tr = tf2.paragraphs[0].runs[0]; tr.font.size = PPt(10); tr.font.color.rgb = GREY
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    foot = s2.shapes.add_textbox(PInches(0.6), PInches(6.85), prs.slide_width - PInches(1.2), PInches(0.5))
    ff = foot.text_frame; ff.word_wrap = True
    ff.text = ("Indicative only; weeks are elapsed time from go-ahead. Pace depends on translator and clinician "
               "availability. Each language goes live as it is completed — value starts at ~week 8, not month 5.")
    fr = ff.paragraphs[0].runs[0]; fr.font.size = PPt(10.5); fr.font.italic = True; fr.font.color.rgb = GREY

    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    cmd = sys.argv[1]
    if cmd == "md2docx":
        md2docx(sys.argv[2], sys.argv[3], sys.argv[4])
    elif cmd == "budget-slide":
        budget_slide(sys.argv[2])
    elif cmd == "language-budget-slide":
        language_budget_slide(sys.argv[2])
